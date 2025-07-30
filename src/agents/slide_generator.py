"""
slide_generator.py — Generates presentation slides from summary points

Uses python-pptx to create a PowerPoint file with a title slide and
bullet-point content slides, saving to the output directory.
"""

import os
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv

load_dotenv()

class SlideGenerator:
    def __init__(self, output_dir = "output/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok = True)

    def create_slide_deck(self, title, bullet_points):
        prs = Presentation()

        # Add title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Add bullet slides, 5 points per slide
        for i in range(0, len(bullet_points)):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]

            title_shape.text = f"{title} (part {i//5 + 1})"
            for bullet in bullet_points[i:i+5]:
                p = content_shape.text_frame.add_paragraph()
                p.text = bullet

        pptx_path = os.path.join(self.output_dir, "generated_presentation.pptx")
        prs.save(pptx_path)
        print(f"\n Presentation saved to: {pptx_path}")

if __name__ == "__main__":
    sample_title = "AI-Based Slide Generation"
    sample_points = [
        "The creation of AI has completely changed the way in how people create presentations.",
        "This project will essentially use the automation that AI provides in order to convert text into slides in a super efficient manner.",
        "The extracted content will then be summarized utilizing numerous different AI techniques."
    ]
    slide_generator = SlideGenerator()
    slide_generator.create_slide_deck(sample_title, sample_points)
